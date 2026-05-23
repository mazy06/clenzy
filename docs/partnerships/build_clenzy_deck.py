#!/usr/bin/env python3
"""
Génère un template PowerPoint pour les rendez-vous partenariat Clenzy.

Style aligné sur la palette Clenzy :
- Primary : #6B8A9A (bleu-gris)
- Accent  : #4A9B8E (vert teal)
- Warm    : #D4A574 (sable)
- Texte   : #2D3748 (slate dark)

Structure 10 slides selon le guide de partnership.
Placeholders [PARTENAIRE], [DATE], [VOLUME], [PRENOM] etc. à remplacer
manuellement à chaque RDV — gardent la cohérence template.

Usage :
    python3 build_clenzy_deck.py
    # → clenzy-partnership-deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ─── Palette Clenzy ────────────────────────────────────────────────────────

PRIMARY   = RGBColor(0x6B, 0x8A, 0x9A)  # bleu-gris
ACCENT    = RGBColor(0x4A, 0x9B, 0x8E)  # vert teal
WARM      = RGBColor(0xD4, 0xA5, 0x74)  # sable
DANGER    = RGBColor(0xC9, 0x7A, 0x7A)
NEUTRAL   = RGBColor(0x8A, 0x83, 0x78)
TEXT      = RGBColor(0x2D, 0x37, 0x48)
TEXT_MUTED= RGBColor(0x71, 0x80, 0x96)
BG_LIGHT  = RGBColor(0xF7, 0xFA, 0xFC)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BORDER    = RGBColor(0xE2, 0xE8, 0xF0)


# ─── Helpers ─────────────────────────────────────────────────────────────


def set_slide_bg(slide, color):
    """Background uniforme pour un slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             font="Helvetica Neue", anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    """Ajoute une textbox avec style standardisé."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    """Rectangle aplat avec fond + bordure optionnelle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    return shape


def add_accent_bar(slide, left, top, width=Inches(0.04), height=Inches(0.4), color=ACCENT):
    """Petite barre verticale accent."""
    return add_rect(slide, left, top, width, height, color)


def add_footer(slide, slide_num, total=10):
    """Footer minimal avec numéro de slide."""
    add_text(slide, Inches(0.4), Inches(7.0), Inches(4), Inches(0.3),
             "Clenzy · Partnership Pitch · v1.0", size=8, color=TEXT_MUTED)
    add_text(slide, Inches(9.5), Inches(7.0), Inches(0.5), Inches(0.3),
             f"{slide_num} / {total}", size=8, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


def add_logo_block(slide, left=Inches(0.4), top=Inches(0.35)):
    """Logo Clenzy textuel + accent bar."""
    add_accent_bar(slide, left, top + Inches(0.05), width=Inches(0.04), height=Inches(0.3), color=ACCENT)
    add_text(slide, left + Inches(0.12), top, Inches(2), Inches(0.4),
             "Clenzy", size=18, bold=True, color=PRIMARY)


# ─── Slide builders ──────────────────────────────────────────────────────


SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)


def build_slide_1_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, WHITE)

    # Big logo centered
    add_text(slide, Inches(0), Inches(1.8), SLIDE_W, Inches(0.8),
             "Clenzy", size=54, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

    # Accent line centered
    add_rect(slide, Inches(4.2), Inches(2.65), Inches(1.6), Inches(0.06), ACCENT)

    # Ribbon "Partnership Proposal"
    add_rect(slide, Inches(3.5), Inches(2.95), Inches(3), Inches(0.4), ACCENT)
    add_text(slide, Inches(3.5), Inches(3.02), Inches(3), Inches(0.35),
             "PARTNERSHIP PROPOSAL", size=10, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    # Subtitle: partenaire name
    add_text(slide, Inches(0), Inches(3.8), SLIDE_W, Inches(0.7),
             "Clenzy × [PARTENAIRE]", size=28, bold=True, color=TEXT,
             align=PP_ALIGN.CENTER)

    # Description
    add_text(slide, Inches(1), Inches(4.7), Inches(8), Inches(0.6),
             "Le PMS du corridor euro-méditerranéen rencontre [Provider]",
             size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # Date + author placeholder
    add_text(slide, Inches(0), Inches(6.5), SLIDE_W, Inches(0.4),
             "[DATE] · [AUTEUR] · Clenzy",
             size=9, color=NEUTRAL, align=PP_ALIGN.CENTER)


def build_slide_2_who_clenzy(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_logo_block(slide)

    # Title
    add_text(slide, Inches(0.4), Inches(1), Inches(9.2), Inches(0.6),
             "Qui est Clenzy ?", size=24, bold=True, color=TEXT)
    add_rect(slide, Inches(0.4), Inches(1.6), Inches(0.6), Inches(0.04), ACCENT)

    # Definition phrase
    add_text(slide, Inches(0.4), Inches(2.0), Inches(9.2), Inches(0.9),
             "Clenzy est un PMS (Property Management System) SaaS multi-tenant "
             "pour la location courte durée, spécialisé sur le corridor "
             "France ↔ Maroc ↔ Arabie Saoudite — un marché tri-régional encore "
             "sous-équipé par les PMS occidentaux.",
             size=14, color=TEXT, line_spacing=1.4)

    # Stats blocks (3 columns)
    stats = [
        ("3", "marchés couverts", "FR · MA · KSA · multi-devises EUR/MAD/SAR"),
        ("17", "catégories d'intégrations", "Signature · Pricing · KYC · Channels · IoT…"),
        ("50+", "fournisseurs branchés", "Pattern Strategy SOLID · API-first"),
    ]
    for i, (n, label, sub) in enumerate(stats):
        x = Inches(0.4 + i * 3.1)
        add_rect(slide, x, Inches(3.5), Inches(2.9), Inches(2.2), BG_LIGHT, BORDER)
        add_accent_bar(slide, x, Inches(3.5), width=Inches(2.9), height=Inches(0.06))
        add_text(slide, x + Inches(0.2), Inches(3.7), Inches(2.5), Inches(0.8),
                 n, size=42, bold=True, color=PRIMARY)
        add_text(slide, x + Inches(0.2), Inches(4.6), Inches(2.5), Inches(0.4),
                 label, size=11, bold=True, color=TEXT)
        add_text(slide, x + Inches(0.2), Inches(5.0), Inches(2.5), Inches(0.7),
                 sub, size=9, color=TEXT_MUTED, line_spacing=1.3)

    # Tagline
    add_text(slide, Inches(0.4), Inches(6.2), Inches(9.2), Inches(0.4),
             "« Un seul PMS pour gérer vos propriétés à Paris, Casablanca et Riyad. »",
             size=12, color=ACCENT, align=PP_ALIGN.CENTER)

    add_footer(slide, 2)


def build_slide_3_market(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_logo_block(slide)

    add_text(slide, Inches(0.4), Inches(1), Inches(9.2), Inches(0.6),
             "Un marché unique sur l'échiquier", size=24, bold=True, color=TEXT)
    add_rect(slide, Inches(0.4), Inches(1.6), Inches(0.6), Inches(0.04), ACCENT)

    add_text(slide, Inches(0.4), Inches(2.0), Inches(9.2), Inches(0.6),
             "France · Maroc · Arabie Saoudite — un corridor STR en pleine "
             "croissance, mal adressé par les PMS US/UK existants.",
             size=12, color=TEXT, line_spacing=1.4)

    # Three market columns
    markets = [
        ("FR", "France", "~800k logements STR",
         "Marché mature\nObligation fiche police\nMulti-OTAs Airbnb/Booking",
         "EUR · Pennylane · Yousign · Chekin · Stripe"),
        ("MA", "Maroc", "~50k logements STR",
         "Marché en croissance rapide\nNumérisation accélérée\nDevise locale MAD",
         "MAD · CMI · Police MA · WhatsApp · Arabe"),
        ("KSA", "Arabie Saoudite", "~30k+ croissance Vision 2030",
         "Marché en explosion (Vision 2030)\nTourisme religieux + business\nDevise locale SAR",
         "SAR · PayTabs · Absher · WhatsApp · Arabe RTL"),
    ]
    for i, (code, name, size_estim, points, stack) in enumerate(markets):
        x = Inches(0.4 + i * 3.1)
        add_rect(slide, x, Inches(2.9), Inches(2.9), Inches(4.0), WHITE, BORDER)
        add_rect(slide, x, Inches(2.9), Inches(2.9), Inches(0.55), PRIMARY)
        add_text(slide, x + Inches(0.2), Inches(2.97), Inches(2.5), Inches(0.4),
                 f"{code} — {name}", size=14, bold=True, color=WHITE)
        add_text(slide, x + Inches(0.2), Inches(3.65), Inches(2.5), Inches(0.3),
                 size_estim, size=10, bold=True, color=ACCENT)
        add_text(slide, x + Inches(0.2), Inches(4.05), Inches(2.5), Inches(1.5),
                 points, size=10, color=TEXT, line_spacing=1.35)
        # Stack technique petit
        add_rect(slide, x + Inches(0.15), Inches(5.85), Inches(2.6), Inches(0.85), BG_LIGHT)
        add_text(slide, x + Inches(0.25), Inches(5.92), Inches(0.6), Inches(0.2),
                 "STACK", size=7, bold=True, color=NEUTRAL)
        add_text(slide, x + Inches(0.25), Inches(6.15), Inches(2.45), Inches(0.55),
                 stack, size=8, color=TEXT_MUTED, line_spacing=1.3)

    add_footer(slide, 3)


def build_slide_4_product(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_logo_block(slide)

    add_text(slide, Inches(0.4), Inches(1), Inches(9.2), Inches(0.6),
             "Produit Clenzy — vue d'ensemble", size=24, bold=True, color=TEXT)
    add_rect(slide, Inches(0.4), Inches(1.6), Inches(0.6), Inches(0.04), ACCENT)

    add_text(slide, Inches(0.4), Inches(2.0), Inches(9.2), Inches(0.5),
             "Trois composants intégrés, conçus pour fonctionner ensemble.",
             size=12, color=TEXT_MUTED)

    products = [
        ("PMS", "Property Management System",
         "Calendrier multi-propriétés\nRéservations + iCal sync\nPricing dynamique\nFacturation NF\nInterventions + ménage",
         "Web app responsive"),
        ("BE", "Booking Engine",
         "Widget JS embarquable\nMulti-langue (FR/EN/AR/RTL)\nPanier multi-séjours\nKYC + paiement\nMulti-devises",
         "SDK pour site host"),
        ("Ops", "Operations Mobile",
         "Check-in/out\nMénage + inventaire\nIntervention technique\nNoise monitoring\nKey management",
         "PWA mobile-first"),
    ]
    for i, (badge, name, features, footer_label) in enumerate(products):
        x = Inches(0.4 + i * 3.1)
        add_rect(slide, x, Inches(2.7), Inches(2.9), Inches(4.0), BG_LIGHT, BORDER)
        # Badge colored top
        badge_color = ACCENT if i == 0 else (PRIMARY if i == 1 else WARM)
        add_rect(slide, x, Inches(2.7), Inches(0.7), Inches(0.7), badge_color)
        add_text(slide, x, Inches(2.83), Inches(0.7), Inches(0.4),
                 badge, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Name
        add_text(slide, x + Inches(0.85), Inches(2.85), Inches(2.0), Inches(0.4),
                 name, size=12, bold=True, color=TEXT)
        # Features
        add_text(slide, x + Inches(0.2), Inches(3.7), Inches(2.6), Inches(2.4),
                 features, size=10, color=TEXT, line_spacing=1.4)
        # Footer
        add_text(slide, x + Inches(0.2), Inches(6.3), Inches(2.6), Inches(0.3),
                 footer_label, size=9, color=TEXT_MUTED)

    add_footer(slide, 4)


def build_slide_5_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_logo_block(slide)

    add_text(slide, Inches(0.4), Inches(1), Inches(9.2), Inches(0.6),
             "Stack technique &amp; intégrations", size=24, bold=True, color=TEXT)
    add_rect(slide, Inches(0.4), Inches(1.6), Inches(0.6), Inches(0.04), ACCENT)

    add_text(slide, Inches(0.4), Inches(2.0), Inches(9.2), Inches(0.5),
             "Architecture moderne, API-first, conçue pour les intégrations.",
             size=12, color=TEXT_MUTED)

    stack_rows = [
        ("Backend",     "Java 21 · Spring Boot 3.2 · JPA · Liquibase · Kafka · Redis"),
        ("Frontend",    "React 18 · TypeScript · MUI · Vite · RTK Query · i18next (FR/EN/AR)"),
        ("Database",    "PostgreSQL 16 · multi-tenant (Hibernate @Filter) · AES-256-GCM"),
        ("Auth",        "Keycloak 24 · OAuth2 · multi-realm (PMS + booking engine)"),
        ("Paiements",   "Stripe (EU) · PayTabs (KSA) · CMI (MA) · Stripe Connect + SEPA payouts"),
        ("Sécurité",    "RGPD · @PreAuthorize partout · validation ownership · NDA-ready"),
    ]
    for i, (label, content) in enumerate(stack_rows):
        y = Inches(2.7 + i * 0.6)
        # Pastille label
        add_rect(slide, Inches(0.4), y, Inches(1.8), Inches(0.45), ACCENT)
        add_text(slide, Inches(0.4), y + Inches(0.08), Inches(1.8), Inches(0.3),
                 label.upper(), size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        # Content
        add_text(slide, Inches(2.4), y + Inches(0.08), Inches(7.2), Inches(0.35),
                 content, size=11, color=TEXT)

    add_text(slide, Inches(0.4), Inches(6.4), Inches(9.2), Inches(0.4),
             "50+ intégrations déjà branchées — Strategy + Registry pattern, ajout d'un provider = 1 classe.",
             size=10, color=ACCENT, align=PP_ALIGN.CENTER)

    add_footer(slide, 5)


def build_slide_6_why_partnership(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_logo_block(slide)

    add_text(slide, Inches(0.4), Inches(1), Inches(9.2), Inches(0.6),
             "Pourquoi ce partenariat", size=24, bold=True, color=TEXT)
    add_rect(slide, Inches(0.4), Inches(1.6), Inches(0.6), Inches(0.04), ACCENT)

    add_text(slide, Inches(0.4), Inches(2.0), Inches(9.2), Inches(0.5),
             "[À personnaliser pour chaque partenaire] — synergie business identifiée.",
             size=12, color=TEXT_MUTED)

    reasons = [
        ("Pour [PARTENAIRE]",
         "▸ Canal de distribution vers le corridor FR/MA/KSA, encore peu adressé\n"
         "▸ Audience qualifiée property managers STR (5-100 logements)\n"
         "▸ Co-marketing : webinaires, études de cas, mentions site Clenzy\n"
         "▸ Volume prévisible : [VOLUME] propriétés en S2 2026"),
        ("Pour Clenzy",
         "▸ Badge partenaire officiel = crédibilité commerciale\n"
         "▸ Tarifs préférentiels à proposer à nos hosts\n"
         "▸ Listing dans le marketplace [PARTENAIRE]\n"
         "▸ Accès à la roadmap produit + support prioritaire"),
    ]
    for i, (title, content) in enumerate(reasons):
        x = Inches(0.4 + i * 4.6)
        add_rect(slide, x, Inches(2.8), Inches(4.4), Inches(3.6), BG_LIGHT, BORDER)
        add_rect(slide, x, Inches(2.8), Inches(4.4), Inches(0.5),
                 ACCENT if i == 0 else PRIMARY)
        add_text(slide, x + Inches(0.2), Inches(2.88), Inches(4), Inches(0.4),
                 title, size=12, bold=True, color=WHITE)
        add_text(slide, x + Inches(0.25), Inches(3.55), Inches(4.05), Inches(2.7),
                 content, size=10.5, color=TEXT, line_spacing=1.5)

    add_text(slide, Inches(0.4), Inches(6.7), Inches(9.2), Inches(0.3),
             "Synergie évidente, pas de chevauchement produit.",
             size=10, color=ACCENT, align=PP_ALIGN.CENTER)

    add_footer(slide, 6)


def build_slide_7_volume(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_logo_block(slide)

    add_text(slide, Inches(0.4), Inches(1), Inches(9.2), Inches(0.6),
             "Volume prévisible &amp; projections", size=24, bold=True, color=TEXT)
    add_rect(slide, Inches(0.4), Inches(1.6), Inches(0.6), Inches(0.04), ACCENT)

    add_text(slide, Inches(0.4), Inches(2.0), Inches(9.2), Inches(0.5),
             "Cibles internes Clenzy sur les 12-18 mois suivants — à actualiser au moment du pitch.",
             size=12, color=TEXT_MUTED)

    quarters = [
        ("Q3 2026", "Phase pilote",   "~30 propriétés",  "~5 organisations"),
        ("Q4 2026", "Lancement",       "~80 propriétés",  "~15 organisations"),
        ("Q1 2027", "Expansion FR",    "~150 propriétés", "~30 organisations"),
        ("Q2 2027", "Ouverture MA/KSA","~250 propriétés", "~50 organisations"),
    ]
    # Header row
    add_rect(slide, Inches(0.4), Inches(2.8), Inches(9.2), Inches(0.5), PRIMARY)
    headers = ["TRIMESTRE", "PHASE", "PROPRIÉTÉS", "ORGANISATIONS"]
    col_widths = [Inches(1.8), Inches(2.7), Inches(2.3), Inches(2.4)]
    x_offset = Inches(0.4)
    for i, h in enumerate(headers):
        add_text(slide, x_offset, Inches(2.88), col_widths[i], Inches(0.3),
                 h, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x_offset += col_widths[i]
    # Rows
    for ri, row in enumerate(quarters):
        y = Inches(3.4 + ri * 0.55)
        bg = WHITE if ri % 2 == 0 else BG_LIGHT
        add_rect(slide, Inches(0.4), y, Inches(9.2), Inches(0.5), bg, BORDER)
        x_offset = Inches(0.4)
        for ci, val in enumerate(row):
            color = ACCENT if ci == 0 else TEXT
            bold = (ci == 0)
            add_text(slide, x_offset, y + Inches(0.1), col_widths[ci], Inches(0.3),
                     val, size=11, bold=bold, color=color, align=PP_ALIGN.CENTER)
            x_offset += col_widths[ci]

    # Bottom info box
    add_rect(slide, Inches(0.4), Inches(6.0), Inches(9.2), Inches(0.7), BG_LIGHT)
    add_accent_bar(slide, Inches(0.4), Inches(6.0), width=Inches(0.06), height=Inches(0.7), color=WARM)
    add_text(slide, Inches(0.6), Inches(6.08), Inches(9), Inches(0.3),
             "Volume estimé impactant pour [PARTENAIRE]", size=10, bold=True, color=TEXT)
    add_text(slide, Inches(0.6), Inches(6.38), Inches(9), Inches(0.3),
             "Sur la base d'un taux de pénétration estimé à [X]% du parc Clenzy → ~[Y] clients référés.",
             size=9, color=TEXT_MUTED)

    add_footer(slide, 7)


def build_slide_8_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_logo_block(slide)

    add_text(slide, Inches(0.4), Inches(1), Inches(9.2), Inches(0.6),
             "Roadmap produit Clenzy", size=24, bold=True, color=TEXT)
    add_rect(slide, Inches(0.4), Inches(1.6), Inches(0.6), Inches(0.04), ACCENT)

    add_text(slide, Inches(0.4), Inches(2.0), Inches(9.2), Inches(0.5),
             "Trois phases sur les 18 mois suivants — orientation produit + partnerships.",
             size=12, color=TEXT_MUTED)

    phases = [
        ("Phase 1 — Foundations", "Maintenant → Q3 2026", ACCENT,
         "• Booking engine multi-langues + RTL\n"
         "• Paiements Stripe + PayTabs + CMI opérationnels\n"
         "• KYC Sumsub natif booking\n"
         "• Intégrations P1 signées"),
        ("Phase 2 — Scale", "Q4 2026 → Q1 2027", PRIMARY,
         "• Marketplace activités cross-sell\n"
         "• Open Banking PIS pour payouts auto\n"
         "• Wise Business pour payouts MA/KSA\n"
         "• Partnerships P2 + co-marketing"),
        ("Phase 3 — Expansion", "Q2 2027+", WARM,
         "• Compliance MA (DGSN) + KSA (Absher)\n"
         "• IoT bundle hardware Clenzy\n"
         "• IA / yield management propriétaire\n"
         "• Marketplace activités complet"),
    ]
    for i, (title, dates, color, content) in enumerate(phases):
        x = Inches(0.4 + i * 3.1)
        add_rect(slide, x, Inches(2.8), Inches(2.9), Inches(3.8), BG_LIGHT, BORDER)
        add_rect(slide, x, Inches(2.8), Inches(2.9), Inches(0.55), color)
        add_text(slide, x + Inches(0.15), Inches(2.88), Inches(2.7), Inches(0.4),
                 title, size=11, bold=True, color=WHITE)
        add_text(slide, x + Inches(0.2), Inches(3.5), Inches(2.6), Inches(0.3),
                 dates, size=9, bold=True, color=color)
        add_text(slide, x + Inches(0.2), Inches(3.85), Inches(2.6), Inches(2.5),
                 content, size=9.5, color=TEXT, line_spacing=1.5)

    add_footer(slide, 8)


def build_slide_9_ask(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_logo_block(slide)

    add_text(slide, Inches(0.4), Inches(1), Inches(9.2), Inches(0.6),
             "Ce que nous demandons", size=24, bold=True, color=TEXT)
    add_rect(slide, Inches(0.4), Inches(1.6), Inches(0.6), Inches(0.04), ACCENT)

    add_text(slide, Inches(0.4), Inches(2.0), Inches(9.2), Inches(0.5),
             "[À adapter au programme partenaires du partenaire] — demande claire et structurée.",
             size=12, color=TEXT_MUTED)

    asks = [
        ("1", "Statut officiel partenaire",
         "Inscription au programme Tech / Solution / Channel partner avec accès aux ressources et au badge marketing."),
        ("2", "Tarif préférentiel partenaire",
         "Discount [X]% à proposer aux clients Clenzy + commission referral sur les abonnements générés."),
        ("3", "Listing marketplace officiel",
         "Présence dans votre PMS Marketplace / Integration Directory avec page Clenzy dédiée."),
        ("4", "Co-marketing structuré",
         "Webinaire conjoint, étude de cas, mention dans la newsletter partenaire, post LinkedIn."),
    ]
    for i, (num, title, desc) in enumerate(asks):
        y = Inches(2.8 + i * 0.85)
        # Numéro circulaire
        add_rect(slide, Inches(0.4), y + Inches(0.05), Inches(0.5), Inches(0.5), ACCENT)
        add_text(slide, Inches(0.4), y + Inches(0.13), Inches(0.5), Inches(0.4),
                 num, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Title
        add_text(slide, Inches(1.1), y + Inches(0.03), Inches(8.5), Inches(0.35),
                 title, size=13, bold=True, color=TEXT)
        # Description
        add_text(slide, Inches(1.1), y + Inches(0.35), Inches(8.5), Inches(0.5),
                 desc, size=10, color=TEXT_MUTED, line_spacing=1.4)

    add_footer(slide, 9)


def build_slide_10_next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_logo_block(slide)

    add_text(slide, Inches(0.4), Inches(1), Inches(9.2), Inches(0.6),
             "Prochaines étapes", size=24, bold=True, color=TEXT)
    add_rect(slide, Inches(0.4), Inches(1.6), Inches(0.6), Inches(0.04), ACCENT)

    # Timeline horizontale
    steps = [
        ("Aujourd'hui",   "Pitch call",          ACCENT),
        ("Sem +1",        "Demo croisée",        PRIMARY),
        ("Sem +2",        "MoU partnership",     WARM),
        ("Sem +4",        "Listing + co-marketing", ACCENT),
    ]
    y_line = Inches(3.2)
    line_left = Inches(1.0)
    line_right = Inches(9.0)
    # Ligne horizontale
    add_rect(slide, line_left, y_line + Inches(0.05), line_right - line_left, Inches(0.04), BORDER)

    step_w = (line_right - line_left) / (len(steps) - 1)
    for i, (date, action, color) in enumerate(steps):
        x_center = line_left + step_w * i
        # Cercle
        circle_size = Inches(0.3)
        add_rect(slide, x_center - circle_size/2, y_line - Inches(0.05), circle_size, circle_size, color)
        # Date au-dessus
        add_text(slide, x_center - Inches(1), y_line - Inches(0.6), Inches(2), Inches(0.3),
                 date, size=9, bold=True, color=color, align=PP_ALIGN.CENTER)
        # Action en-dessous
        add_text(slide, x_center - Inches(1), y_line + Inches(0.5), Inches(2), Inches(0.5),
                 action, size=11, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

    # Contact card
    add_rect(slide, Inches(2), Inches(4.7), Inches(6), Inches(1.8), BG_LIGHT, BORDER)
    add_accent_bar(slide, Inches(2), Inches(4.7), width=Inches(0.06), height=Inches(1.8), color=ACCENT)

    add_text(slide, Inches(2.3), Inches(4.85), Inches(5.5), Inches(0.4),
             "Restons en contact", size=14, bold=True, color=TEXT)
    add_text(slide, Inches(2.3), Inches(5.25), Inches(5.5), Inches(0.4),
             "[VOTRE PRÉNOM NOM] — [VOTRE TITRE], Clenzy",
             size=11, color=TEXT)
    add_text(slide, Inches(2.3), Inches(5.6), Inches(5.5), Inches(0.4),
             "📧  [votre.email]@clenzy.fr", size=10, color=TEXT_MUTED)
    add_text(slide, Inches(2.3), Inches(5.9), Inches(5.5), Inches(0.4),
             "📞  +33 [votre.téléphone]", size=10, color=TEXT_MUTED)
    add_text(slide, Inches(2.3), Inches(6.2), Inches(5.5), Inches(0.4),
             "🌐  https://clenzy.fr", size=10, color=TEXT_MUTED)

    # Thank you
    add_text(slide, Inches(0), Inches(6.9), SLIDE_W, Inches(0.4),
             "Merci pour votre temps.",
             size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ─── Main ────────────────────────────────────────────────────────────────


def build_deck(output_path: Path):
    prs = Presentation()
    # 16:9 format
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_1_cover(prs)
    build_slide_2_who_clenzy(prs)
    build_slide_3_market(prs)
    build_slide_4_product(prs)
    build_slide_5_stack(prs)
    build_slide_6_why_partnership(prs)
    build_slide_7_volume(prs)
    build_slide_8_roadmap(prs)
    build_slide_9_ask(prs)
    build_slide_10_next_steps(prs)

    prs.save(str(output_path))
    print(f"✓ Deck généré : {output_path}")
    print(f"  Slides : {len(prs.slides)}")
    print(f"  Format : 16:9 ({prs.slide_width.inches}\" × {prs.slide_height.inches}\")")


if __name__ == "__main__":
    output = Path(__file__).parent / "clenzy-partnership-deck.pptx"
    build_deck(output)
